import re
from pathlib import Path


def _desktop_directory():
    desktop = Path.home() / "Desktop"
    desktop.mkdir(parents=True, exist_ok=True)
    return desktop.resolve()


def crea_presentazione(titolo, diapositive, nome_file=None):
    from pptx import Presentation
    from pptx.util import Inches, Pt

    title = str(titolo or "Presentazione").strip()[:140]
    slides = list(diapositive or [])
    if not slides:
        return {"successo": False, "messaggio": "Servono almeno una diapositiva e i relativi contenuti."}
    safe = re.sub(r"[^\w .-]+", "_", str(nome_file or title), flags=re.UNICODE).strip(" .") or "Presentazione"
    if not safe.lower().endswith(".pptx"):
        safe += ".pptx"
    desktop = _desktop_directory()
    target = (desktop / safe).resolve()
    if target.parent != desktop:
        return {"successo": False, "messaggio": "Nome file non valido."}
    deck = Presentation()
    deck.slide_width, deck.slide_height = Inches(13.333), Inches(7.5)
    cover = deck.slides.add_slide(deck.slide_layouts[0])
    cover.shapes.title.text = title
    cover.placeholders[1].text = "Creata da JARVIS"
    for row in slides[:40]:
        item = row if isinstance(row, dict) else {"titolo": str(row), "contenuto": []}
        slide = deck.slides.add_slide(deck.slide_layouts[1])
        slide.shapes.title.text = str(item.get("titolo") or "Sezione")[:160]
        body = slide.placeholders[1].text_frame
        body.clear()
        content = item.get("contenuto") or item.get("punti") or []
        if isinstance(content, str):
            content = [line.strip() for line in content.splitlines() if line.strip()]
        for index, point in enumerate(list(content)[:12]):
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = str(point)[:700]
            paragraph.font.size = Pt(24)
    deck.save(target)
    return {"successo": True, "messaggio": f"Presentazione creata sul Desktop: {target.name}.",
            "dati": {"percorso": str(target), "diapositive": len(deck.slides)}}
